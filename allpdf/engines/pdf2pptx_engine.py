# allpdf/engines/pdf2pptx_engine.py
"""PDF to PPTX engine — renders each PDF page as a full-slide image."""
import os
import time
import tempfile
from pathlib import Path

import fitz
from pptx import Presentation
from pptx.util import Inches

from allpdf.engines.base import ConversionEngine
from allpdf.models import ConversionResult, ConversionStatus, FileFormat


class Pdf2PptxEngine(ConversionEngine):
    """Convert PDF to PPTX by rendering each page as a high-DPI image.

    Each PDF page becomes one slide. Uses 16:9 widescreen format.
    """

    name = "pdf2pptx"
    input_format = FileFormat.PDF
    output_format = FileFormat.PPTX

    def convert(self, input_path: str, output_path: str, **options) -> ConversionResult:
        start = time.time()
        dpi = options.get("dpi", 200)

        if not os.path.exists(input_path):
            return ConversionResult(
                input_path=Path(input_path),
                output_path=Path(output_path),
                input_format=FileFormat.PDF,
                output_format=FileFormat.PPTX,
                status=ConversionStatus.FAILED,
                engine_used=self.name,
                duration_seconds=time.time() - start,
                error_message=f"Input file not found: {input_path}",
            )

        try:
            doc = fitz.open(input_path)
            prs = Presentation()
            prs.slide_width = Inches(13.333)
            prs.slide_height = Inches(7.5)

            pages = options.get("pages", list(range(doc.page_count)))

            with tempfile.TemporaryDirectory() as tmpdir:
                for page_num in pages:
                    page = doc[page_num]
                    mat = fitz.Matrix(dpi / 72, dpi / 72)
                    pix = page.get_pixmap(matrix=mat)
                    img_path = os.path.join(tmpdir, f"page_{page_num}.png")
                    pix.save(img_path)

                    slide_layout = prs.slide_layouts[6]  # blank layout
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.add_picture(
                        img_path,
                        Inches(0), Inches(0),
                        width=prs.slide_width,
                        height=prs.slide_height,
                    )

            doc.close()
            prs.save(output_path)

            duration = time.time() - start

            return ConversionResult(
                input_path=Path(input_path),
                output_path=Path(output_path),
                input_format=FileFormat.PDF,
                output_format=FileFormat.PPTX,
                status=ConversionStatus.SUCCESS,
                engine_used=self.name,
                duration_seconds=duration,
            )

        except Exception as e:
            duration = time.time() - start
            return ConversionResult(
                input_path=Path(input_path),
                output_path=Path(output_path),
                input_format=FileFormat.PDF,
                output_format=FileFormat.PPTX,
                status=ConversionStatus.FAILED,
                engine_used=self.name,
                duration_seconds=duration,
                error_message=str(e),
            )

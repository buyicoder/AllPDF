# allpdf/quality/checker.py
"""Quality checker — validates conversion results across multiple dimensions."""
import os
from pathlib import Path

from allpdf.models import (
    FileFormat, FileInfo, QualityCheckItem, QualityGrade, QualityReport,
)
from allpdf.engines.pymupdf_ops import PyMuPDFOps


class QualityChecker:
    """Runs quality checks on conversion results.

    Checks: page count, file openability, image count, text extractability.
    """

    def __init__(self):
        self._ops = PyMuPDFOps()

    def run(
        self,
        input_path: str,
        output_path: str,
        input_format: FileFormat,
        output_format: FileFormat,
    ) -> QualityReport:
        """Run all quality checks and return a report."""
        checks: list[QualityCheckItem] = []

        source_info = self._ops.analyze(input_path)
        if source_info is None:
            source_info = FileInfo(
                path=Path(input_path), format=input_format,
                page_count=0, image_count=0, table_count=0,
                is_scanned=False, file_size_bytes=0,
            )

        result_info = self._analyze_result(output_path, output_format)

        checks.append(self.check_page_count(source_info, result_info, input_format, output_format))
        checks.append(self.check_openable(output_path, output_format))

        if source_info.image_count > 0:
            checks.append(self.check_image_count(source_info, result_info))

        checks.append(self.check_text_extractable(output_path, output_format))

        return QualityReport.from_checks(checks)

    def check_page_count(
        self,
        source: FileInfo,
        result: FileInfo,
        input_fmt: FileFormat,
        output_fmt: FileFormat,
    ) -> QualityCheckItem:
        """Check page counts. Office->PDF: exact. PDF->Office: +/-1 tolerance."""
        diff = abs(source.page_count - result.page_count)
        if input_fmt.is_office_format() and output_fmt == FileFormat.PDF:
            if diff == 0:
                return QualityCheckItem(name="page_count", grade=QualityGrade.GREEN,
                    detail=f"{source.page_count} pages, exact match")
            else:
                return QualityCheckItem(name="page_count", grade=QualityGrade.YELLOW,
                    detail=f"Expected {source.page_count}, got {result.page_count}")
        else:
            if diff <= 1:
                return QualityCheckItem(name="page_count", grade=QualityGrade.GREEN,
                    detail=f"Source {source.page_count}, result {result.page_count} (within tolerance)")
            else:
                return QualityCheckItem(name="page_count", grade=QualityGrade.YELLOW,
                    detail=f"Page count mismatch: {source.page_count} vs {result.page_count}")

    def check_openable(self, path: str, fmt: FileFormat) -> QualityCheckItem:
        """Verify the output file can be opened by its format library."""
        if not os.path.exists(path):
            return QualityCheckItem(name="file_openable", grade=QualityGrade.RED,
                detail="Output file does not exist")

        try:
            if fmt == FileFormat.DOCX:
                from docx import Document
                Document(path)
            elif fmt == FileFormat.XLSX:
                from openpyxl import load_workbook
                wb = load_workbook(path, read_only=True)
                wb.close()
            elif fmt == FileFormat.PPTX:
                from pptx import Presentation
                Presentation(path)
            elif fmt == FileFormat.PDF:
                import fitz
                doc = fitz.open(path)
                doc.close()
            return QualityCheckItem(name="file_openable", grade=QualityGrade.GREEN,
                detail="File opens successfully")
        except Exception as e:
            return QualityCheckItem(name="file_openable", grade=QualityGrade.RED,
                detail=f"Cannot open: {e}")

    def check_image_count(self, source: FileInfo, result: FileInfo) -> QualityCheckItem:
        """Compare embedded image counts."""
        if source.image_count == result.image_count:
            return QualityCheckItem(name="image_count", grade=QualityGrade.GREEN,
                detail=f"{source.image_count} images, matches")
        else:
            return QualityCheckItem(name="image_count", grade=QualityGrade.YELLOW,
                detail=f"Expected {source.image_count} images, got {result.image_count}")

    def check_text_extractable(self, path: str, fmt: FileFormat) -> QualityCheckItem:
        """Verify text can be extracted from the result."""
        try:
            if fmt == FileFormat.DOCX:
                from docx import Document
                doc = Document(path)
                text = " ".join(p.text for p in doc.paragraphs)
            elif fmt == FileFormat.XLSX:
                from openpyxl import load_workbook
                wb = load_workbook(path, read_only=True)
                text = ""
                for ws in wb.worksheets:
                    for row in ws.iter_rows(max_row=5):
                        for cell in row:
                            if cell.value:
                                text += str(cell.value) + " "
                wb.close()
            elif fmt == FileFormat.PPTX:
                from pptx import Presentation
                prs = Presentation(path)
                text = ""
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if shape.has_text_frame:
                            text += shape.text_frame.text + " "
            elif fmt == FileFormat.PDF:
                import fitz
                doc = fitz.open(path)
                text = ""
                for page in doc:
                    text += page.get_text()
                doc.close()

            if text.strip():
                return QualityCheckItem(name="text_extractable", grade=QualityGrade.GREEN,
                    detail=f"Extracted {len(text)} chars")
            else:
                return QualityCheckItem(name="text_extractable", grade=QualityGrade.YELLOW,
                    detail="No text could be extracted")
        except Exception as e:
            return QualityCheckItem(name="text_extractable", grade=QualityGrade.RED,
                detail=f"Text extraction failed: {e}")

    def _analyze_result(self, path: str, fmt: FileFormat) -> FileInfo:
        """Analyze the result file using format-specific tools."""
        file_size = os.path.getsize(path) if os.path.exists(path) else 0

        if fmt == FileFormat.PDF:
            info = self._ops.analyze(path)
            if info:
                return info

        page_count = 0
        image_count = 0

        try:
            if fmt == FileFormat.DOCX:
                from allpdf.utils import count_docx_pages
                page_count = count_docx_pages(path)
            elif fmt == FileFormat.XLSX:
                from allpdf.utils import count_xlsx_sheets
                page_count = count_xlsx_sheets(path)
            elif fmt == FileFormat.PPTX:
                from allpdf.utils import count_pptx_slides
                page_count = count_pptx_slides(path)
        except Exception:
            pass

        return FileInfo(
            path=Path(path), format=fmt,
            page_count=page_count, image_count=image_count,
            table_count=0, is_scanned=False, file_size_bytes=file_size,
        )
